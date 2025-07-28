"""
Export Service for Deep Research application.

This service handles exporting research reports to various formats:
- Markdown (raw and formatted)
- PDF generation using ReportLab
- PowerPoint presentations using python-pptx with custom templates

Includes template management and Azure Storage integration for file hosting.
"""

import asyncio
import os
import tempfile
import uuid
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any

import structlog
from jinja2 import Environment, FileSystemLoader, Template
import markdown
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.colors import HexColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
import aiofiles

from app.core.azure_config import AzureServiceManager
from app.models.schemas import ResearchReport, ResearchSection, ExportFormat


logger = structlog.get_logger(__name__)


class ExportService:
    """
    Service for exporting research reports to various formats.
    
    Supports:
    - Markdown export with custom formatting
    - PDF generation with professional styling
    - Word Document (DOCX) generation
    - PowerPoint presentations with template support
    - HTML export
    - JSON export
    - Azure Storage integration for file hosting
    """
    
    def __init__(self, azure_manager: AzureServiceManager):
        """
        Initialize Export Service.
        
        Args:
            azure_manager: Azure service manager instance
        """
        self.azure_manager = azure_manager
        
        # Template configuration
        self.templates_dir = Path(__file__).parent.parent / "templates"
        self.templates_dir.mkdir(exist_ok=True)
        
        # Initialize Jinja2 environment
        self.jinja_env = Environment(
            loader=FileSystemLoader(str(self.templates_dir)),
            autoescape=True
        )
        
        # Export directory for temporary files
        self.export_dir = Path(tempfile.gettempdir()) / "deep_research_exports"
        self.export_dir.mkdir(exist_ok=True)
        
        # Template configurations
        self.pptx_templates = {
            "default": "default_template.pptx",
            "business": "business_template.pptx",
            "academic": "academic_template.pptx",
            "executive": "executive_template.pptx"
        }
        
        # Ensure templates exist
        asyncio.create_task(self._ensure_templates_exist())
    
    async def export_markdown(
        self,
        report: ResearchReport,
        export_id: str,
        include_metadata: bool = True
    ) -> str:
        """
        Export research report as Markdown.
        
        Args:
            report: Research report to export
            export_id: Export task identifier
            include_metadata: Whether to include report metadata
            
        Returns:
            Path to the generated Markdown file
        """
        try:
            logger.info("Exporting report as Markdown", export_id=export_id, task_id=report.task_id)
            
            # Generate Markdown content
            markdown_content = await self._generate_markdown_content(report, include_metadata)
            
            # Save to file
            file_path = self.export_dir / f"report_{export_id}.md"
            
            async with aiofiles.open(file_path, 'w', encoding='utf-8') as f:
                await f.write(markdown_content)
            
            logger.info("Markdown export completed", export_id=export_id, file_path=str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            logger.error("Markdown export failed", export_id=export_id, error=str(e), exc_info=True)
            raise
    
    async def export_pdf(
        self,
        report: ResearchReport,
        export_id: str,
        include_sources: bool = True,
        include_metadata: bool = True
    ) -> str:
        """
        Export research report as PDF.
        
        Args:
            report: Research report to export
            export_id: Export task identifier
            include_sources: Whether to include source citations
            include_metadata: Whether to include report metadata
            
        Returns:
            Path to the generated PDF file
        """
        try:
            logger.info("Exporting report as PDF", export_id=export_id, task_id=report.task_id)
            
            # Generate PDF using ReportLab
            file_path = self.export_dir / f"report_{export_id}.pdf"
            
            await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: self._generate_pdf_with_reportlab(
                    report, str(file_path), include_sources, include_metadata
                )
            )
            
            logger.info("PDF export completed", export_id=export_id, file_path=str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            logger.error("PDF export failed", export_id=export_id, error=str(e), exc_info=True)
            raise
    
    async def export_pptx(
        self,
        report: ResearchReport,
        export_id: str,
        template_name: Optional[str] = None,
        custom_branding: Optional[Dict[str, str]] = None
    ) -> str:
        """
        Export research report as PowerPoint presentation.
        
        Args:
            report: Research report to export
            export_id: Export task identifier
            template_name: PPTX template to use
            custom_branding: Custom branding options
            
        Returns:
            Path to the generated PPTX file
        """
        try:
            logger.info(
                "Exporting report as PPTX",
                export_id=export_id,
                task_id=report.task_id,
                template=template_name
            )
            
            # Load template
            template_path = await self._get_pptx_template(template_name or "default")
            
            # Create presentation from template
            prs = Presentation(template_path)
            
            # Generate slides based on report content
            await self._populate_pptx_slides(prs, report, custom_branding)
            
            # Save presentation
            file_path = self.export_dir / f"report_{export_id}.pptx"
            
            await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: prs.save(str(file_path))
            )
            
            logger.info("PPTX export completed", export_id=export_id, file_path=str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            logger.error("PPTX export failed", export_id=export_id, error=str(e), exc_info=True)
            raise
    
    async def export_docx(
        self,
        report: ResearchReport,
        export_id: str,
        include_sources: bool = True,
        include_metadata: bool = True,
        include_table_of_contents: bool = True,
        include_page_numbers: bool = True
    ) -> str:
        """
        Export research report as Word document.
        
        Args:
            report: Research report to export
            export_id: Export task identifier
            include_sources: Whether to include source citations
            include_metadata: Whether to include report metadata
            include_table_of_contents: Whether to include table of contents
            include_page_numbers: Whether to include page numbers
            
        Returns:
            Path to the generated DOCX file
        """
        try:
            logger.info("Exporting report as DOCX", export_id=export_id, task_id=report.task_id)
            
            # Generate DOCX using python-docx
            file_path = self.export_dir / f"report_{export_id}.docx"
            
            await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: self._generate_docx_with_python_docx(
                    report, str(file_path), include_sources, include_metadata,
                    include_table_of_contents, include_page_numbers
                )
            )
            
            logger.info("DOCX export completed", export_id=export_id, file_path=str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            logger.error("DOCX export failed", export_id=export_id, error=str(e), exc_info=True)
            raise

    async def export_html(
        self,
        report: ResearchReport,
        export_id: str,
        include_sources: bool = True,
        include_metadata: bool = True,
        custom_css: Optional[str] = None
    ) -> str:
        """
        Export research report as HTML.
        
        Args:
            report: Research report to export
            export_id: Export task identifier
            include_sources: Whether to include source citations
            include_metadata: Whether to include report metadata
            custom_css: Custom CSS for styling
            
        Returns:
            Path to the generated HTML file
        """
        try:
            logger.info("Exporting report as HTML", export_id=export_id, task_id=report.task_id)
            
            # Generate HTML content
            file_path = self.export_dir / f"report_{export_id}.html"
            
            html_content = await self._generate_html_content(
                report, include_sources, include_metadata, custom_css
            )
            
            async with aiofiles.open(file_path, 'w', encoding='utf-8') as f:
                await f.write(html_content)
            
            logger.info("HTML export completed", export_id=export_id, file_path=str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            logger.error("HTML export failed", export_id=export_id, error=str(e), exc_info=True)
            raise

    async def export_json(
        self,
        report: ResearchReport,
        export_id: str,
        include_raw_data: bool = True
    ) -> str:
        """
        Export research report as JSON.
        
        Args:
            report: Research report to export
            export_id: Export task identifier
            include_raw_data: Whether to include raw report data
            
        Returns:
            Path to the generated JSON file
        """
        try:
            logger.info("Exporting report as JSON", export_id=export_id, task_id=report.task_id)
            
            # Generate JSON content
            file_path = self.export_dir / f"report_{export_id}.json"
            
            json_content = report.model_dump() if include_raw_data else {
                "task_id": report.task_id,
                "title": report.title,
                "executive_summary": report.executive_summary,
                "sections": [
                    {
                        "title": section.title,
                        "content": section.content,
                        "word_count": section.word_count,
                        "confidence_score": section.confidence_score
                    }
                    for section in report.sections
                ],
                "conclusions": report.conclusions,
                "word_count": report.word_count,
                "reading_time_minutes": report.reading_time_minutes,
                "created_at": report.created_at.isoformat()
            }
            
            async with aiofiles.open(file_path, 'w', encoding='utf-8') as f:
                await f.write(json.dumps(json_content, indent=2, default=str))
            
            logger.info("JSON export completed", export_id=export_id, file_path=str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            logger.error("JSON export failed", export_id=export_id, error=str(e), exc_info=True)
            raise
    
    async def _generate_markdown_content(
        self,
        report: ResearchReport,
        include_metadata: bool
    ) -> str:
        """Generate formatted Markdown content for the report."""
        lines = []
        
        # Title
        lines.append(f"# {report.title}")
        lines.append("")
        
        # Metadata
        if include_metadata:
            lines.append("## Report Information")
            lines.append("")
            lines.append(f"- **Generated**: {report.created_at.strftime('%Y-%m-%d %H:%M:%S UTC')}")
            lines.append(f"- **Task ID**: `{report.task_id}`")
            lines.append(f"- **Word Count**: {report.word_count:,}")
            lines.append(f"- **Reading Time**: {report.reading_time_minutes} minutes")
            if report.metadata:
                for key, value in report.metadata.items():
                    lines.append(f"- **{key.replace('_', ' ').title()}**: {value}")
            lines.append("")
        
        # Executive Summary
        lines.append("## Executive Summary")
        lines.append("")
        lines.append(report.executive_summary)
        lines.append("")
        
        # Sections
        for section in report.sections:
            lines.append(f"## {section.title}")
            lines.append("")
            lines.append(section.content)
            lines.append("")
            
            # Add sources if available
            if section.sources:
                lines.append("### Sources")
                lines.append("")
                for i, source in enumerate(section.sources, 1):
                    lines.append(f"{i}. [{source.title}]({source.url})")
                    if source.snippet:
                        lines.append(f"   _{source.snippet}_")
                lines.append("")
        
        # Conclusions
        if report.conclusions:
            lines.append("## Conclusions")
            lines.append("")
            lines.append(report.conclusions)
            lines.append("")
        
        # All Sources
        if report.sources and include_metadata:
            lines.append("## References")
            lines.append("")
            for i, source in enumerate(report.sources, 1):
                lines.append(f"{i}. [{source.title}]({source.url})")
                if source.snippet:
                    lines.append(f"   _{source.snippet}_")
                if source.published_date:
                    lines.append(f"   Published: {source.published_date.strftime('%Y-%m-%d')}")
                lines.append("")
        
        return "\n".join(lines)
    
    def _generate_pdf_with_reportlab(
        self,
        report: ResearchReport,
        file_path: str,
        include_sources: bool,
        include_metadata: bool
    ) -> None:
        """Generate PDF using ReportLab."""
        doc = SimpleDocTemplate(file_path, pagesize=A4)
        story = []
        
        # Get styles
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=HexColor('#1a365d'),
            alignment=1  # Center alignment
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            spaceBefore=20,
            textColor=HexColor('#2d3748')
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=12,
            leading=14
        )
        
        # Title
        story.append(Paragraph(report.title, title_style))
        story.append(Spacer(1, 12))
        
        # Metadata
        if include_metadata:
            if report.created_at:
                story.append(Paragraph(f"<b>Generated:</b> {report.created_at.strftime('%Y-%m-%d %H:%M')}", body_style))
            if report.task_id:
                story.append(Paragraph(f"<b>Task ID:</b> {report.task_id}", body_style))
            story.append(Spacer(1, 20))
        
        # Summary
        if report.summary:
            story.append(Paragraph("Executive Summary", heading_style))
            story.append(Paragraph(report.summary, body_style))
            story.append(Spacer(1, 20))
        
        # Sections
        for section in report.sections:
            story.append(Paragraph(section.title, heading_style))
            
            # Clean up content for PDF
            content = section.content.replace('\n\n', '<br/><br/>')
            content = content.replace('\n', ' ')
            
            story.append(Paragraph(content, body_style))
            story.append(Spacer(1, 15))
        
        # Sources
        if include_sources and report.sources:
            story.append(PageBreak())
            story.append(Paragraph("Sources", heading_style))
            
            for i, source in enumerate(report.sources, 1):
                source_text = f"<b>[{i}]</b> {source.title}"
                if source.url:
                    source_text += f"<br/><i>{source.url}</i>"
                if source.author:
                    source_text += f"<br/>Author: {source.author}"
                if source.published_date:
                    source_text += f"<br/>Published: {source.published_date.strftime('%Y-%m-%d')}"
                
                story.append(Paragraph(source_text, body_style))
                story.append(Spacer(1, 10))
        
        # Build PDF
        doc.build(story)
    
    async def _populate_pptx_slides(
        self,
        prs: Presentation,
        report: ResearchReport,
        custom_branding: Optional[Dict[str, str]]
    ) -> None:
        """Populate PowerPoint slides with report content."""
        # Clear existing slides (keep only title slide)
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, 0, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
        
        # Title slide
        title_slide = prs.slides[0]
        await self._populate_title_slide(title_slide, report, custom_branding)
        
        # Executive Summary slide
        summary_layout = prs.slide_layouts[1]  # Title and Content layout
        summary_slide = prs.slides.add_slide(summary_layout)
        await self._populate_summary_slide(summary_slide, report)
        
        # Section slides
        for section in report.sections:
            section_layout = prs.slide_layouts[1]  # Title and Content layout
            section_slide = prs.slides.add_slide(section_layout)
            await self._populate_section_slide(section_slide, section)
        
        # Key Findings slide (if applicable)
        findings_layout = prs.slide_layouts[1]
        findings_slide = prs.slides.add_slide(findings_layout)
        await self._populate_findings_slide(findings_slide, report)
        
        # Sources slide
        if report.sources:
            sources_layout = prs.slide_layouts[1]
            sources_slide = prs.slides.add_slide(sources_layout)
            await self._populate_sources_slide(sources_slide, report)
    
    async def _populate_title_slide(
        self,
        slide,
        report: ResearchReport,
        custom_branding: Optional[Dict[str, str]]
    ) -> None:
        """Populate the title slide."""
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = report.title
        
        subtitle_text = f"Deep Research Report\n"
        subtitle_text += f"Generated: {report.created_at.strftime('%B %d, %Y')}\n"
        subtitle_text += f"Reading Time: {report.reading_time_minutes} minutes"
        
        if custom_branding and "company" in custom_branding:
            subtitle_text += f"\n\nPrepared by: {custom_branding['company']}"
        
        subtitle.text = subtitle_text
    
    async def _populate_summary_slide(self, slide, report: ResearchReport) -> None:
        """Populate the executive summary slide."""
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Summary"
        content.text = report.executive_summary
    
    async def _populate_section_slide(self, slide, section: ResearchSection) -> None:
        """Populate a section slide."""
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = section.title
        
        # Clean markdown formatting for PowerPoint
        clean_content = section.content.replace('**', '').replace('*', '').replace('#', '')
        
        # Truncate if too long
        if len(clean_content) > 500:
            clean_content = clean_content[:497] + "..."
        
        content.text = clean_content
    
    async def _populate_findings_slide(self, slide, report: ResearchReport) -> None:
        """Populate key findings slide."""
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Key Findings"
        
        # Extract key points from sections
        key_points = []
        for section in report.sections:
            if "finding" in section.title.lower() or "key" in section.title.lower():
                # Extract bullet points or first few sentences
                lines = section.content.split('\n')
                for line in lines:
                    if line.strip().startswith(('-', '*', '1.', '2.', '3.')):
                        key_points.append(line.strip())
                    elif len(line.strip()) > 20 and len(key_points) < 5:
                        key_points.append(f"• {line.strip()[:100]}...")
        
        if not key_points:
            key_points = [f"• {report.conclusions[:100]}..."]
        
        content.text = '\n'.join(key_points[:6])  # Limit to 6 points
    
    async def _populate_sources_slide(self, slide, report: ResearchReport) -> None:
        """Populate sources slide."""
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Sources"
        
        source_list = []
        for i, source in enumerate(report.sources[:8], 1):  # Limit to 8 sources
            source_list.append(f"{i}. {source.title}")
            if source.domain:
                source_list.append(f"   {source.domain}")
        
        content.text = '\n'.join(source_list)
    
    async def _get_pptx_template(self, template_name: str) -> str:
        """Get path to PPTX template file."""
        template_file = self.pptx_templates.get(template_name, self.pptx_templates["default"])
        template_path = self.templates_dir / template_file
        
        if not template_path.exists():
            # Create a basic template if it doesn't exist
            await self._create_default_pptx_template(template_path)
        
        return str(template_path)
    
    async def _create_default_pptx_template(self, template_path: Path) -> None:
        """Create a default PPTX template."""
        try:
            # Create a basic presentation template
            prs = Presentation()
            
            # Customize the default slide master
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Deep Research Report"
            subtitle.text = "Professional Research Analysis"
            
            # Save the template
            await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: prs.save(str(template_path))
            )
            
            logger.info("Created default PPTX template", template_path=str(template_path))
            
        except Exception as e:
            logger.error("Failed to create PPTX template", error=str(e))
    
    async def _ensure_templates_exist(self) -> None:
        """Ensure all required templates exist."""
        try:
            for template_name, template_file in self.pptx_templates.items():
                template_path = self.templates_dir / template_file
                if not template_path.exists():
                    await self._create_default_pptx_template(template_path)
            
            logger.info("All export templates verified")
            
        except Exception as e:
            logger.error("Failed to ensure templates exist", error=str(e))
    
    async def upload_to_azure_storage(self, file_path: str, blob_name: str) -> str:
        """
        Upload exported file to Azure Blob Storage.
        
        Args:
            file_path: Local file path
            blob_name: Blob name in storage
            
        Returns:
            Public URL of the uploaded file
        """
        try:
            blob_client = self.azure_manager.blob_client
            if not blob_client:
                raise ValueError("Azure Blob Storage not configured")
            
            container_name = "exports"
            
            # Upload file
            with open(file_path, 'rb') as data:
                blob_client.get_blob_client(
                    container=container_name,
                    blob=blob_name
                ).upload_blob(data, overwrite=True)
            
            # Generate public URL
            account_url = self.azure_manager.settings.STORAGE_ACCOUNT_URL
            public_url = f"{account_url}/{container_name}/{blob_name}"
            
            logger.info("File uploaded to Azure Storage", blob_name=blob_name, url=public_url)
            
            return public_url
            
        except Exception as e:
            logger.error("Failed to upload to Azure Storage", file_path=file_path, error=str(e))
            raise
    
    def cleanup_export_file(self, file_path: str) -> None:
        """Clean up temporary export file."""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                logger.debug("Export file cleaned up", file_path=file_path)
        except Exception as e:
            logger.warning("Failed to cleanup export file", file_path=file_path, error=str(e))
    
    def _generate_docx_with_python_docx(
        self,
        report: ResearchReport,
        file_path: str,
        include_sources: bool,
        include_metadata: bool,
        include_table_of_contents: bool,
        include_page_numbers: bool
    ) -> None:
        """Generate Word document using python-docx."""
        try:
            # Create document
            doc = Document()
            
            # Configure styles
            styles = doc.styles
            
            # Title style
            title_style = styles.add_style('ReportTitle', WD_STYLE_TYPE.PARAGRAPH)
            title_font = title_style.font
            title_font.name = 'Arial'
            title_font.size = DocxPt(24)
            title_font.bold = True
            title_style.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_style.paragraph_format.space_after = DocxPt(12)
            
            # Heading styles
            heading1_style = styles.add_style('ReportHeading1', WD_STYLE_TYPE.PARAGRAPH)
            heading1_font = heading1_style.font
            heading1_font.name = 'Arial'
            heading1_font.size = DocxPt(18)
            heading1_font.bold = True
            heading1_style.paragraph_format.space_before = DocxPt(12)
            heading1_style.paragraph_format.space_after = DocxPt(6)
            
            heading2_style = styles.add_style('ReportHeading2', WD_STYLE_TYPE.PARAGRAPH)
            heading2_font = heading2_style.font
            heading2_font.name = 'Arial'
            heading2_font.size = DocxPt(14)
            heading2_font.bold = True
            heading2_style.paragraph_format.space_before = DocxPt(10)
            heading2_style.paragraph_format.space_after = DocxPt(4)
            
            # Body style
            body_style = styles.add_style('ReportBody', WD_STYLE_TYPE.PARAGRAPH)
            body_font = body_style.font
            body_font.name = 'Arial'
            body_font.size = DocxPt(11)
            body_style.paragraph_format.space_after = DocxPt(6)
            body_style.paragraph_format.line_spacing = 1.15
            
            # Add title
            title_para = doc.add_paragraph(report.title, style='ReportTitle')
            
            # Add metadata if requested
            if include_metadata:
                doc.add_paragraph('Report Information', style='ReportHeading1')
                
                metadata_para = doc.add_paragraph(style='ReportBody')
                metadata_para.add_run(f"Generated: {report.created_at.strftime('%Y-%m-%d %H:%M:%S UTC')}\n")
                metadata_para.add_run(f"Task ID: {report.task_id}\n")
                metadata_para.add_run(f"Word Count: {report.word_count:,}\n")
                metadata_para.add_run(f"Reading Time: {report.reading_time_minutes} minutes")
            
            # Add executive summary
            doc.add_paragraph('Executive Summary', style='ReportHeading1')
            doc.add_paragraph(report.executive_summary, style='ReportBody')
            
            # Add sections
            for section in report.sections:
                doc.add_paragraph(section.title, style='ReportHeading2')
                
                # Split content into paragraphs and process
                content_paragraphs = section.content.split('\n\n')
                for paragraph_text in content_paragraphs:
                    if paragraph_text.strip():
                        # Simple markdown processing
                        processed_text = paragraph_text.replace('**', '').replace('*', '')
                        doc.add_paragraph(processed_text, style='ReportBody')
                
                # Add sources if requested
                if include_sources and section.sources:
                    sources_para = doc.add_paragraph(style='ReportBody')
                    sources_para.add_run("Sources:\n").bold = True
                    for i, source in enumerate(section.sources, 1):
                        sources_para.add_run(f"{i}. {source.title} - {source.url}\n")
            
            # Add conclusions
            if report.conclusions:
                doc.add_paragraph('Conclusions', style='ReportHeading1')
                doc.add_paragraph(report.conclusions, style='ReportBody')
            
            # Add sources section if requested
            if include_sources and report.sources:
                doc.add_paragraph('References', style='ReportHeading1')
                for i, source in enumerate(report.sources, 1):
                    source_para = doc.add_paragraph(style='ReportBody')
                    source_para.add_run(f"{i}. ").bold = True
                    source_para.add_run(f"{source.title}\n")
                    source_para.add_run(f"   {source.url}\n")
                    if source.snippet:
                        source_para.add_run(f"   {source.snippet[:100]}...")
            
            # Save document
            doc.save(file_path)
            
        except Exception as e:
            logger.error("Failed to generate DOCX with python-docx", error=str(e))
            raise
    
    async def _generate_html_content(
        self,
        report: ResearchReport,
        include_sources: bool,
        include_metadata: bool,
        custom_css: Optional[str] = None
    ) -> str:
        """Generate HTML content for the report."""
        try:
            # Default CSS
            default_css = """
            <style>
                body { font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }
                h1 { color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }
                h2 { color: #34495e; border-bottom: 1px solid #bdc3c7; padding-bottom: 5px; }
                h3 { color: #7f8c8d; }
                .metadata { background-color: #f8f9fa; padding: 15px; border-radius: 5px; margin: 20px 0; }
                .section { margin: 30px 0; }
                .sources { background-color: #f0f0f0; padding: 10px; border-radius: 5px; margin: 10px 0; }
                .source-item { margin: 5px 0; }
                .source-url { color: #3498db; text-decoration: none; }
                .source-url:hover { text-decoration: underline; }
            </style>
            """
            
            css = custom_css if custom_css else default_css
            
            html_lines = [
                "<!DOCTYPE html>",
                "<html>",
                "<head>",
                "<meta charset='utf-8'>",
                f"<title>{report.title}</title>",
                css,
                "</head>",
                "<body>",
                f"<h1>{report.title}</h1>"
            ]
            
            # Add metadata if requested
            if include_metadata:
                html_lines.extend([
                    "<div class='metadata'>",
                    "<h2>Report Information</h2>",
                    f"<p><strong>Generated:</strong> {report.created_at.strftime('%Y-%m-%d %H:%M:%S UTC')}</p>",
                    f"<p><strong>Task ID:</strong> {report.task_id}</p>",
                    f"<p><strong>Word Count:</strong> {report.word_count:,}</p>",
                    f"<p><strong>Reading Time:</strong> {report.reading_time_minutes} minutes</p>",
                    "</div>"
                ])
            
            # Add executive summary
            html_lines.extend([
                "<h2>Executive Summary</h2>",
                f"<p>{report.executive_summary}</p>"
            ])
            
            # Add sections
            for section in report.sections:
                html_lines.extend([
                    "<div class='section'>",
                    f"<h2>{section.title}</h2>"
                ])
                
                # Convert markdown to HTML
                content_html = markdown.markdown(section.content)
                html_lines.append(content_html)
                
                # Add sources if requested
                if include_sources and section.sources:
                    html_lines.extend([
                        "<div class='sources'>",
                        "<h3>Sources:</h3>"
                    ])
                    for i, source in enumerate(section.sources, 1):
                        html_lines.append(
                            f"<div class='source-item'>{i}. "
                            f"<a href='{source.url}' class='source-url' target='_blank'>{source.title}</a>"
                            f"</div>"
                        )
                    html_lines.append("</div>")
                
                html_lines.append("</div>")
            
            # Add conclusions
            if report.conclusions:
                html_lines.extend([
                    "<h2>Conclusions</h2>",
                    f"<p>{report.conclusions}</p>"
                ])
            
            # Add references section if requested
            if include_sources and report.sources:
                html_lines.extend([
                    "<h2>References</h2>",
                    "<div class='sources'>"
                ])
                for i, source in enumerate(report.sources, 1):
                    html_lines.extend([
                        f"<div class='source-item'>",
                        f"<strong>{i}.</strong> ",
                        f"<a href='{source.url}' class='source-url' target='_blank'>{source.title}</a>",
                        f"<br><small>{source.snippet[:100] if source.snippet else ''}...</small>",
                        f"</div>"
                    ])
                html_lines.append("</div>")
            
            html_lines.extend([
                "</body>",
                "</html>"
            ])
            
            return "\n".join(html_lines)
            
        except Exception as e:
            logger.error("Failed to generate HTML content", error=str(e))
            raise

    async def create_custom_powerpoint(
        self,
        slides_data: Dict[str, Any],
        topic: str,
        template_name: str = "business"
    ) -> str:
        """
        Create a custom PowerPoint presentation from structured slide data.
        
        Args:
            slides_data: JSON structure containing slides with titles and content
            topic: Research topic for the presentation title
            template_name: PowerPoint template to use
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            # Get template path
            template_path = self.templates_dir / self.pptx_templates.get(template_name, "business_template.pptx")
            
            # Create presentation from template or new if template doesn't exist
            if template_path.exists():
                prs = Presentation(str(template_path))
                logger.info(f"Using PowerPoint template: {template_path}")
            else:
                prs = Presentation()
                logger.warning(f"Template not found, creating new presentation: {template_path}")
            
            # Remove existing slides except the first one (title slide)
            slide_count = len(prs.slides)
            for i in range(slide_count - 1, 0, -1):  # Remove from last to first (except slide 0)
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
            
            # Update title slide if exists
            if len(prs.slides) > 0:
                title_slide = prs.slides[0]
                if title_slide.shapes.title:
                    title_slide.shapes.title.text = f"Research Report: {topic}"
                if len(title_slide.placeholders) > 1:
                    title_slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            
            # Add slides from the structured data
            slides = slides_data.get("slides", [])
            
            for slide_info in slides:
                slide_title = slide_info.get("title", "Untitled Slide")
                slide_content = slide_info.get("content", [])
                
                # Add new slide with content layout
                slide_layout = prs.slide_layouts[1]  # Typically "Title and Content" layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set slide title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                
                # Add content to the slide
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    
                    # Handle different content types
                    if isinstance(slide_content, list):
                        # Regular bullet points
                        if content_placeholder.has_text_frame:
                            text_frame = content_placeholder.text_frame
                            text_frame.clear()
                            
                            for i, bullet_point in enumerate(slide_content):
                                if isinstance(bullet_point, str):
                                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                                    p.text = bullet_point
                                    p.level = 0
                    
                    elif isinstance(slide_content, dict):
                        # Special handling for structured content like SWOT analysis
                        if content_placeholder.has_text_frame:
                            text_frame = content_placeholder.text_frame
                            text_frame.clear()
                            
                            paragraph_added = False
                            for category, items in slide_content.items():
                                # Add category header
                                p = text_frame.paragraphs[0] if not paragraph_added else text_frame.add_paragraph()
                                p.text = f"{category}:"
                                p.level = 0
                                paragraph_added = True
                                
                                # Add items under category
                                if isinstance(items, list):
                                    for item in items:
                                        p = text_frame.add_paragraph()
                                        p.text = str(item)
                                        p.level = 1
                    
                    else:
                        # Single content item
                        if content_placeholder.has_text_frame:
                            content_placeholder.text = str(slide_content)
                
                logger.info(f"Added slide: {slide_title}")
            
            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).rstrip()[:30]
            filename = f"custom_pptx_{safe_topic}_{timestamp}.pptx"
            output_path = self.export_dir / filename
            
            # Save the presentation
            prs.save(str(output_path))
            
            logger.info(f"Custom PowerPoint created successfully: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error("Failed to create custom PowerPoint", error=str(e), exc_info=True)
            raise
