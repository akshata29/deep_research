"""
Convert API endpoints for document format conversion.

Handles conversion between different document formats:
- Markdown to PDF using pdfkit (wkhtmltopdf)
- Markdown to HTML
- Markdown to DOCX
"""

import asyncio
import os
import tempfile
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional

import structlog
from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import FileResponse
from pydantic import BaseModel
import aiofiles
import markdown
import pdfkit
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


router = APIRouter()
logger = structlog.get_logger(__name__)


class MarkdownConvertRequest(BaseModel):
    """Request model for markdown conversion."""
    markdown_content: str
    title: str
    css_style: Optional[str] = None


@router.post("/markdown-to-pdf")
async def convert_markdown_to_pdf(request: MarkdownConvertRequest):
    """
    Convert Markdown content to PDF.
    
    Args:
        request: Markdown conversion request
        
    Returns:
        PDF file response
    """
    try:
        logger.info("Converting markdown to PDF", title=request.title)
        
        # Create temporary directory for this conversion
        temp_dir = Path(tempfile.gettempdir()) / f"conversion_{uuid.uuid4()}"
        temp_dir.mkdir(exist_ok=True)
        
        # Convert markdown to HTML first
        html_content = await _markdown_to_html(
            request.markdown_content, 
            request.title, 
            request.css_style
        )
        
        # Save HTML to temporary file
        html_file = temp_dir / "content.html"
        async with aiofiles.open(html_file, 'w', encoding='utf-8') as f:
            await f.write(html_content)
        
        # Convert HTML to PDF using pdfkit
        pdf_file = temp_dir / f"{request.title}.pdf"
        
        await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: _html_to_pdf_pdfkit(str(html_file), str(pdf_file))
        )
        
        logger.info("PDF conversion completed", file_path=str(pdf_file))
        
        # Return the PDF file
        return FileResponse(
            path=str(pdf_file),
            filename=f"{request.title}.pdf",
            media_type="application/pdf",
            background=lambda: asyncio.create_task(_cleanup_temp_dir(temp_dir))
        )
        
    except Exception as e:
        logger.error("PDF conversion failed", error=str(e), exc_info=True)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to convert to PDF: {str(e)}"
        )


@router.post("/markdown-to-html")
async def convert_markdown_to_html(request: MarkdownConvertRequest):
    """
    Convert Markdown content to HTML.
    
    Args:
        request: Markdown conversion request
        
    Returns:
        HTML file response
    """
    try:
        logger.info("Converting markdown to HTML", title=request.title)
        
        # Create temporary directory for this conversion
        temp_dir = Path(tempfile.gettempdir()) / f"conversion_{uuid.uuid4()}"
        temp_dir.mkdir(exist_ok=True)
        
        # Convert markdown to HTML
        html_content = await _markdown_to_html(
            request.markdown_content, 
            request.title, 
            request.css_style
        )
        
        # Save HTML to temporary file
        html_file = temp_dir / f"{request.title}.html"
        async with aiofiles.open(html_file, 'w', encoding='utf-8') as f:
            await f.write(html_content)
        
        logger.info("HTML conversion completed", file_path=str(html_file))
        
        # Return the HTML file
        return FileResponse(
            path=str(html_file),
            filename=f"{request.title}.html",
            media_type="text/html",
            background=lambda: asyncio.create_task(_cleanup_temp_dir(temp_dir))
        )
        
    except Exception as e:
        logger.error("HTML conversion failed", error=str(e), exc_info=True)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to convert to HTML: {str(e)}"
        )


@router.post("/markdown-to-docx")
async def convert_markdown_to_docx(request: MarkdownConvertRequest):
    """
    Convert Markdown content to DOCX (Word document).
    
    Args:
        request: Markdown conversion request
        
    Returns:
        DOCX file response
    """
    try:
        logger.info("Converting markdown to DOCX", title=request.title)
        
        # Create temporary directory for this conversion
        temp_dir = Path(tempfile.gettempdir()) / f"conversion_{uuid.uuid4()}"
        temp_dir.mkdir(exist_ok=True)
        
        # Convert markdown to DOCX
        docx_file = temp_dir / f"{request.title}.docx"
        
        await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: _markdown_to_docx(request.markdown_content, request.title, str(docx_file))
        )
        
        logger.info("DOCX conversion completed", file_path=str(docx_file))
        
        # Return the DOCX file
        return FileResponse(
            path=str(docx_file),
            filename=f"{request.title}.docx",
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            background=lambda: asyncio.create_task(_cleanup_temp_dir(temp_dir))
        )
        
    except Exception as e:
        logger.error("DOCX conversion failed", error=str(e), exc_info=True)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to convert to DOCX: {str(e)}"
        )


@router.post("/markdown-to-pptx")
async def convert_markdown_to_pptx(request: MarkdownConvertRequest):
    """
    Convert Markdown content to PPTX (PowerPoint presentation).
    
    Args:
        request: Markdown conversion request
        
    Returns:
        PPTX file response
    """
    try:
        logger.info("Converting markdown to PPTX", title=request.title)
        
        # Create temporary directory for this conversion
        temp_dir = Path(tempfile.gettempdir()) / f"conversion_{uuid.uuid4()}"
        temp_dir.mkdir(exist_ok=True)
        
        # Convert markdown to PPTX
        pptx_file = temp_dir / f"{request.title}.pptx"
        
        await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: _markdown_to_pptx(request.markdown_content, request.title, str(pptx_file))
        )
        
        logger.info("PPTX conversion completed", file_path=str(pptx_file))
        
        # Return the PPTX file
        return FileResponse(
            path=str(pptx_file),
            filename=f"{request.title}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            background=lambda: asyncio.create_task(_cleanup_temp_dir(temp_dir))
        )
        
    except Exception as e:
        logger.error("PPTX conversion failed", error=str(e), exc_info=True)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to convert to PPTX: {str(e)}"
        )


async def _markdown_to_html(markdown_content: str, title: str, custom_css: Optional[str] = None) -> str:
    """Convert markdown content to styled HTML."""
    
    # Default CSS for professional-looking documents
    default_css = """
    <style>
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            line-height: 1.6;
            margin: 40px auto;
            max-width: 800px;
            padding: 20px;
            color: #333;
            background-color: #fff;
        }
        
        h1 {
            color: #2c3e50;
            border-bottom: 3px solid #3498db;
            padding-bottom: 10px;
            margin-top: 0;
            page-break-after: avoid;
        }
        
        h2 {
            color: #34495e;
            border-bottom: 2px solid #bdc3c7;
            padding-bottom: 5px;
            margin-top: 30px;
            page-break-after: avoid;
        }
        
        h3 {
            color: #7f8c8d;
            margin-top: 25px;
            page-break-after: avoid;
        }
        
        h4, h5, h6 {
            color: #95a5a6;
            margin-top: 20px;
            page-break-after: avoid;
        }
        
        p {
            margin-bottom: 16px;
            text-align: justify;
        }
        
        code {
            background-color: #f8f9fa;
            padding: 2px 6px;
            border-radius: 3px;
            font-family: 'Consolas', 'Monaco', 'Courier New', monospace;
            font-size: 0.9em;
            color: #e74c3c;
        }
        
        pre {
            background-color: #f8f9fa;
            padding: 15px;
            border-radius: 5px;
            overflow-x: auto;
            border-left: 4px solid #3498db;
            margin: 20px 0;
            page-break-inside: avoid;
        }
        
        pre code {
            background: none;
            padding: 0;
            color: #333;
            font-size: 0.85em;
        }
        
        blockquote {
            border-left: 4px solid #3498db;
            margin: 20px 0;
            padding: 10px 20px;
            background-color: #f8f9fa;
            font-style: italic;
            page-break-inside: avoid;
        }
        
        ul, ol {
            margin-bottom: 16px;
            padding-left: 30px;
        }
        
        li {
            margin-bottom: 8px;
        }
        
        table {
            border-collapse: collapse;
            width: 100%;
            margin: 20px 0;
            page-break-inside: avoid;
        }
        
        th, td {
            border: 1px solid #ddd;
            padding: 12px;
            text-align: left;
        }
        
        th {
            background-color: #f2f2f2;
            font-weight: bold;
        }
        
        tr:nth-child(even) {
            background-color: #f9f9f9;
        }
        
        a {
            color: #3498db;
            text-decoration: none;
        }
        
        a:hover {
            text-decoration: underline;
        }
        
        .title-page {
            text-align: center;
            margin-bottom: 40px;
            page-break-after: always;
        }
        
        .footer {
            margin-top: 40px;
            text-align: center;
            font-size: 0.9em;
            color: #7f8c8d;
            border-top: 1px solid #bdc3c7;
            padding-top: 20px;
        }
        
        @media print {
            body {
                margin: 0;
                padding: 20px;
            }
            
            .page-break {
                page-break-before: always;
            }
            
            h1, h2, h3 {
                page-break-after: avoid;
            }
            
            table, pre, blockquote {
                page-break-inside: avoid;
            }
        }
    </style>
    """
    
    # Convert markdown to HTML
    md = markdown.Markdown(extensions=[
        'tables',
        'fenced_code',
        'codehilite',
        'toc',
        'footnotes',
        'attr_list',
        'def_list'
    ])
    
    html_body = md.convert(markdown_content)
    
    # Create complete HTML document
    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    {custom_css if custom_css else default_css}
</head>
<body>
    <div class="title-page">
        <h1>{title}</h1>
        <p>Generated on {datetime.now().strftime('%B %d, %Y')}</p>
    </div>
    
    <div class="content">
        {html_body}
    </div>
    
    <div class="footer">
        <p>Generated by Deep Research - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
    </div>
</body>
</html>"""
    
    return html_content


def _html_to_pdf_pdfkit(html_file_path: str, pdf_file_path: str) -> None:
    """Convert HTML file to PDF using pdfkit (wkhtmltopdf)."""
    try:
        # Configuration options for wkhtmltopdf
        options = {
            'page-size': 'A4',
            'margin-top': '0.75in',
            'margin-right': '0.75in',
            'margin-bottom': '0.75in',
            'margin-left': '0.75in',
            'encoding': "UTF-8",
            'no-outline': None,
            'enable-local-file-access': None,
            'quiet': '',
        }
        
        # Try to find wkhtmltopdf executable on Windows
        possible_paths = [
            r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe",
            r"C:\Program Files (x86)\wkhtmltopdf\bin\wkhtmltopdf.exe",
            "wkhtmltopdf"  # If it's in PATH
        ]
        
        wkhtmltopdf_path = None
        for path in possible_paths:
            if os.path.exists(path) or path == "wkhtmltopdf":
                wkhtmltopdf_path = path
                break
        
        if wkhtmltopdf_path and wkhtmltopdf_path != "wkhtmltopdf":
            config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path)
            pdfkit.from_file(html_file_path, pdf_file_path, options=options, configuration=config)
        else:
            # Try without explicit path (assumes wkhtmltopdf is in PATH)
            pdfkit.from_file(html_file_path, pdf_file_path, options=options)
        
    except Exception as e:
        logger.error("pdfkit conversion failed", error=str(e))
        # Fallback: try without options if the conversion fails
        try:
            pdfkit.from_file(html_file_path, pdf_file_path)
        except Exception as fallback_error:
            logger.error("pdfkit fallback conversion failed", error=str(fallback_error))
            raise Exception(
                "PDF conversion failed. Please ensure wkhtmltopdf is installed. "
                "Download from: https://wkhtmltopdf.org/downloads.html"
            )


def _markdown_to_docx(markdown_content: str, title: str, docx_file_path: str) -> None:
    """Convert markdown content to DOCX using python-docx."""
    try:
        # Create a new document
        doc = Document()
        
        # Add custom styles
        _add_docx_styles(doc)
        
        # Add title
        title_paragraph = doc.add_heading(title, 0)
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add generation timestamp
        timestamp_paragraph = doc.add_paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
        timestamp_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add a page break
        doc.add_page_break()
        
        # Parse markdown content line by line
        lines = markdown_content.split('\n')
        current_list = None
        in_code_block = False
        code_content = []
        
        i = 0
        while i < len(lines):
            line = lines[i].rstrip()
            
            # Handle code blocks
            if line.startswith('```'):
                if in_code_block:
                    # End of code block
                    _add_code_block_to_docx(doc, '\n'.join(code_content))
                    code_content = []
                    in_code_block = False
                else:
                    # Start of code block
                    in_code_block = True
                i += 1
                continue
            
            if in_code_block:
                code_content.append(line)
                i += 1
                continue
            
            # Handle headings
            if line.startswith('#'):
                current_list = None  # Reset list context
                level = len(line) - len(line.lstrip('#'))
                heading_text = line.lstrip('#').strip()
                if heading_text:
                    doc.add_heading(heading_text, level)
            
            # Handle lists
            elif line.strip().startswith(('- ', '* ', '+ ')):
                list_text = line.strip()[2:].strip()
                if current_list != 'bullet':
                    current_list = 'bullet'
                doc.add_paragraph(list_text, style='List Bullet')
            
            elif line.strip() and line.strip()[0].isdigit() and '. ' in line.strip():
                list_text = line.strip().split('. ', 1)[1] if '. ' in line.strip() else line.strip()
                if current_list != 'number':
                    current_list = 'number'
                doc.add_paragraph(list_text, style='List Number')
            
            # Handle blockquotes
            elif line.strip().startswith('>'):
                current_list = None
                quote_text = line.strip()[1:].strip()
                if quote_text:
                    quote_paragraph = doc.add_paragraph(quote_text)
                    quote_paragraph.style = 'Quote'
            
            # Handle regular paragraphs
            elif line.strip():
                current_list = None
                # Process inline formatting
                processed_line = _process_inline_formatting(line.strip())
                if processed_line:
                    paragraph = doc.add_paragraph()
                    _add_formatted_text_to_paragraph(paragraph, processed_line)
            
            # Handle empty lines
            else:
                current_list = None
                if i > 0 and lines[i-1].strip():  # Add spacing only after content
                    doc.add_paragraph()
            
            i += 1
        
        # Add footer
        section = doc.sections[0]
        footer = section.footer
        footer_paragraph = footer.paragraphs[0]
        footer_paragraph.text = f"Generated by Deep Research - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        footer_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Save the document
        doc.save(docx_file_path)
        
    except Exception as e:
        logger.error("Failed to convert markdown to DOCX", error=str(e))
        raise Exception(f"DOCX conversion failed: {str(e)}")


def _add_docx_styles(doc: Document) -> None:
    """Add custom styles to the document."""
    try:
        styles = doc.styles
        
        # Add a custom quote style
        if 'Quote' not in [style.name for style in styles]:
            quote_style = styles.add_style('Quote', WD_STYLE_TYPE.PARAGRAPH)
            quote_style.base_style = styles['Normal']
            quote_paragraph_format = quote_style.paragraph_format
            quote_paragraph_format.left_indent = Inches(0.5)
            quote_paragraph_format.right_indent = Inches(0.5)
            quote_font = quote_style.font
            quote_font.italic = True
            quote_font.color.rgb = None  # Keep default color
        
        # Add a custom code style
        if 'Code Block' not in [style.name for style in styles]:
            code_style = styles.add_style('Code Block', WD_STYLE_TYPE.PARAGRAPH)
            code_style.base_style = styles['Normal']
            code_paragraph_format = code_style.paragraph_format
            code_paragraph_format.left_indent = Inches(0.25)
            code_paragraph_format.right_indent = Inches(0.25)
            code_font = code_style.font
            code_font.name = 'Consolas'
            code_font.size = Inches(0.15)  # Slightly smaller font
            
    except Exception as e:
        logger.warning("Failed to add custom styles", error=str(e))


def _add_code_block_to_docx(doc: Document, code_content: str) -> None:
    """Add a code block to the document."""
    try:
        paragraph = doc.add_paragraph()
        paragraph.style = 'Code Block' if 'Code Block' in [style.name for style in doc.styles] else 'Normal'
        run = paragraph.add_run(code_content)
        run.font.name = 'Consolas'
        
    except Exception as e:
        logger.warning("Failed to add code block", error=str(e))
        # Fallback: add as regular paragraph
        doc.add_paragraph(code_content)


def _process_inline_formatting(text: str) -> str:
    """Process inline markdown formatting."""
    # This is a simple implementation - for production you might want to use a proper markdown parser
    # For now, we'll handle basic formatting
    return text


def _add_formatted_text_to_paragraph(paragraph, text: str) -> None:
    """Add formatted text to a paragraph, handling basic markdown formatting."""
    try:
        # Simple implementation - you can enhance this to handle bold, italic, links, etc.
        # For now, just add the text as-is
        paragraph.add_run(text)
        
    except Exception as e:
        logger.warning("Failed to add formatted text", error=str(e))
        paragraph.add_run(text)


def _markdown_to_pptx(markdown_content: str, title: str, pptx_file_path: str) -> None:
    """Convert markdown content to PPTX using python-pptx."""
    try:
        # Create a new presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        subtitle_placeholder.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Parse markdown content
        lines = markdown_content.split('\n')
        current_slide = None
        current_content = []
        slide_title = ""
        
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            
            # Handle headings (create new slides)
            if line.startswith('#'):
                # Save previous slide if it has content
                if current_slide is not None and current_content:
                    _add_content_to_slide(current_slide, current_content)
                
                # Create new slide
                level = len(line) - len(line.lstrip('#'))
                slide_title = line.lstrip('#').strip()
                
                if slide_title and level <= 2:  # Only create slides for H1 and H2
                    slide_layout = prs.slide_layouts[1]  # Title and Content layout
                    current_slide = prs.slides.add_slide(slide_layout)
                    current_slide.shapes.title.text = slide_title
                    current_content = []
            
            # Handle content
            elif line and current_slide is not None:
                current_content.append(line)
            
            # Handle lists and other content
            elif line:
                if current_slide is None:
                    # Create a content slide if we don't have one
                    slide_layout = prs.slide_layouts[1]
                    current_slide = prs.slides.add_slide(slide_layout)
                    current_slide.shapes.title.text = "Content"
                    current_content = []
                
                current_content.append(line)
            
            i += 1
        
        # Add final slide content
        if current_slide is not None and current_content:
            _add_content_to_slide(current_slide, current_content)
        
        # If no slides were created, create a single content slide
        if len(prs.slides) == 1:  # Only title slide
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Content"
            _add_content_to_slide(slide, markdown_content.split('\n'))
        
        # Save the presentation
        prs.save(pptx_file_path)
        
    except Exception as e:
        logger.error("Failed to convert markdown to PPTX", error=str(e))
        raise Exception(f"PPTX conversion failed: {str(e)}")


def _add_content_to_slide(slide, content_lines):
    """Add content lines to a PowerPoint slide."""
    try:
        # Get the content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if hasattr(shape, 'text') and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder is None:
            # Create a text box if no content placeholder
            left = PptxInches(1)
            top = PptxInches(2)
            width = PptxInches(8)
            height = PptxInches(4)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
        
        # Add content
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, line in enumerate(content_lines):
            line = line.strip()
            if not line:
                continue
                
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Handle different types of content
            if line.startswith(('- ', '* ', '+ ')):
                # Bullet point
                p.text = line[2:].strip()
                p.level = 0
            elif line.startswith('  - ') or line.startswith('    - '):
                # Sub-bullet point
                p.text = line.strip()[2:]
                p.level = 1
            elif line.startswith(('1. ', '2. ', '3. ', '4. ', '5. ')):
                # Numbered list
                p.text = line[3:].strip()
                p.level = 0
            else:
                # Regular text
                p.text = line
                p.level = 0
            
            # Basic formatting
            if p.text:
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(18)
                
    except Exception as e:
        logger.warning("Failed to add content to slide", error=str(e))


async def _cleanup_temp_dir(temp_dir: Path) -> None:
    """Clean up temporary directory after file is served."""
    try:
        # Add a small delay to ensure file is served before cleanup
        await asyncio.sleep(1)
        
        # Remove all files in the directory
        for file_path in temp_dir.iterdir():
            if file_path.is_file():
                file_path.unlink()
        
        # Remove the directory
        temp_dir.rmdir()
        
        logger.debug("Temporary directory cleaned up", path=str(temp_dir))
        
    except Exception as e:
        logger.warning("Failed to cleanup temporary directory", path=str(temp_dir), error=str(e))
